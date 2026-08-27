#!/usr/bin/env python3
"""Fast deterministic checks for generated or edited .pptx files.

Usage:
    python preflight_check.py <pptx_file> [--slides 1,3,9]
    python preflight_check.py <pptx_file> --integrity [--expected-slides N]
    python preflight_check.py <pptx_file> --integrity --geometry [--slides 1,3,9]

Modes:
    default / --geometry  Geometry checks on all slides, or on --slides scope.
    --integrity           Full-file package/relationship/slide-count checks only.
    --integrity --geometry
                          Run both check families in one invocation.

Exit 0 = pass; Exit 1 = issues found.

Design principle: geometry checks flag only text-bearing objects. Decorative
shapes extending beyond the canvas may be intentional bleed and are not treated
as defects here.
"""
from __future__ import annotations

import argparse
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu


def check_integrity(path: Path, expected_slides: int | None = None) -> list[str]:
    """Return structure/package issues for the whole presentation."""
    issues: list[str] = []

    try:
        with zipfile.ZipFile(path, "r") as zf:
            names = zf.namelist()
            seen: set[str] = set()
            for name in names:
                if name in seen:
                    issues.append(f"  ZIP: duplicate entry '{name}'")
                seen.add(name)
            bad_member = zf.testzip()
            if bad_member:
                issues.append(f"  ZIP: CRC failure in '{bad_member}'")
    except zipfile.BadZipFile as exc:
        return [f"  ZIP: corrupt file — {exc}"]

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        return issues + [f"  PARSE: python-pptx cannot open — {exc}"]

    actual = len(prs.slides)
    if expected_slides is not None and actual != expected_slides:
        issues.append(f"  COUNT: expected {expected_slides} slides, got {actual}")

    for idx, slide in enumerate(prs.slides, 1):
        try:
            for rel in slide.part.rels.values():
                if rel.is_external:
                    continue
                try:
                    _ = rel.target_part
                except Exception:
                    issues.append(
                        f"  S{idx}: dangling rel {rel.reltype} -> {rel.target_ref}"
                    )
        except Exception as exc:
            issues.append(f"  S{idx}: relationship scan failed — {exc}")

    return issues


def has_text(shape) -> bool:
    return bool(shape.has_text_frame and shape.text_frame.text.strip())


def iter_shapes(shapes):
    """Yield shapes recursively so grouped text is checked too."""
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def check_bounds(slide, idx, slide_w, slide_h, issues: list[str]) -> None:
    for shape in iter_shapes(slide.shapes):
        if not has_text(shape):
            continue
        try:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            if None in (left, top, width, height):
                continue
            tol = Emu(0.1)
            sample = shape.text_frame.text.strip().replace("\n", " ")[:20]
            if left < -tol or top < -tol:
                issues.append(f"  S{idx}: text '{sample}' has negative position")
            if left + width > slide_w + tol:
                delta = (left + width - slide_w) / 914400
                issues.append(f"  S{idx}: text '{sample}' past right by {delta:.2f}\"")
            if top + height > slide_h + tol:
                delta = (top + height - slide_h) / 914400
                issues.append(f"  S{idx}: text '{sample}' past bottom by {delta:.2f}\"")
        except Exception:
            continue


def check_text_overlaps(slide, idx, issues: list[str]) -> None:
    boxes = []
    for shape in iter_shapes(slide.shapes):
        if not has_text(shape):
            continue
        try:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            if None in (left, top, width, height):
                continue
            sample = shape.text_frame.text.strip().replace("\n", " ")[:25]
            boxes.append((sample, left, top, left + width, top + height, width * height))
        except Exception:
            continue

    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            text1, l1, t1, r1, b1, a1 = boxes[i]
            text2, l2, t2, r2, b2, a2 = boxes[j]
            if l1 < r2 and l2 < r1 and t1 < b2 and t2 < b1:
                overlap_w = min(r1, r2) - max(l1, l2)
                overlap_h = min(b1, b2) - max(t1, t2)
                min_area = min(a1, a2)
                if min_area > 0 and (overlap_w * overlap_h) / min_area > 0.30:
                    issues.append(f"  S{idx}: text overlap '{text1}' <-> '{text2}'")


def check_empty_text_boxes(slide, idx, issues: list[str]) -> None:
    placeholders = {"", "xxx", "todo", "[insert", "lorem"}
    for shape in iter_shapes(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
            continue
        text = shape.text_frame.text.strip()
        if text.lower() in placeholders:
            issues.append(f"  S{idx}: empty/placeholder TEXT_BOX (forgot to fill?)")


def check_fonts(slide, idx, issues: list[str]) -> None:
    for shape in iter_shapes(slide.shapes):
        if not has_text(shape):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip() and run.font.name is None:
                    issues.append(f"  S{idx}: no fontFace on '{run.text[:20]}'")


def check_geometry(path: Path, slide_scope: set[int] | None = None) -> list[str]:
    issues: list[str] = []
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        return [f"  PARSE: python-pptx cannot open — {exc}"]

    slide_w, slide_h = prs.slide_width, prs.slide_height
    for idx, slide in enumerate(prs.slides, 1):
        if slide_scope and idx not in slide_scope:
            continue
        check_bounds(slide, idx, slide_w, slide_h, issues)
        check_text_overlaps(slide, idx, issues)
        check_empty_text_boxes(slide, idx, issues)
        check_fonts(slide, idx, issues)
    return issues


def parse_slide_scope(raw: str | None) -> set[int] | None:
    if not raw:
        return None
    try:
        result = {int(x.strip()) for x in raw.split(",") if x.strip()}
    except ValueError as exc:
        raise argparse.ArgumentTypeError("--slides must be comma-separated integers") from exc
    if any(x < 1 for x in result):
        raise argparse.ArgumentTypeError("slide numbers are 1-based and must be >= 1")
    return result


def main() -> None:
    parser = argparse.ArgumentParser(description="Deterministic preflight checks for .pptx files")
    parser.add_argument("pptx_file", type=Path)
    parser.add_argument("--integrity", action="store_true", help="run full-file package/integrity checks")
    parser.add_argument("--geometry", action="store_true", help="run geometry/text checks; this is the default when no mode flag is given")
    parser.add_argument("--slides", help="1-based comma-separated slide scope for geometry checks, e.g. 1,3,9")
    parser.add_argument("--expected-slides", type=int, help="expected slide count; valid with --integrity")
    args = parser.parse_args()

    if not args.pptx_file.is_file():
        parser.error(f"file not found: {args.pptx_file}")
    if args.pptx_file.suffix.lower() != ".pptx":
        parser.error("pptx_file must end in .pptx")
    if args.expected_slides is not None and not args.integrity:
        parser.error("--expected-slides requires --integrity")

    slide_scope = parse_slide_scope(args.slides)
    if slide_scope and args.integrity and not args.geometry:
        parser.error("--slides applies to geometry checks; add --geometry or omit --integrity")

    run_integrity = args.integrity
    run_geometry = args.geometry or not args.integrity

    issues: list[str] = []
    if run_integrity:
        issues.extend(check_integrity(args.pptx_file, args.expected_slides))
    if run_geometry:
        issues.extend(check_geometry(args.pptx_file, slide_scope))

    if issues:
        print("PREFLIGHT ISSUES:")
        for issue in issues:
            print(issue)
        print(f"\n{len(issues)} issue(s). Fix before visual QA.")
        sys.exit(1)

    modes = []
    if run_integrity:
        modes.append("integrity")
    if run_geometry:
        modes.append("geometry")
    print(f"Preflight: ALL CHECKS PASSED ({' + '.join(modes)})")


if __name__ == "__main__":
    main()
